import json
import logging
from pathlib import Path
from typing import List, Dict, Any, Optional
import fitz  # PyMuPDF
import docx
import openpyxl
import pptx

# Configure logger
logger = logging.getLogger("docmind_ai")

class DocumentService:
    @staticmethod
    def extract_text_from_bytes(file_bytes: bytes, file_name: str) -> Optional[List[Dict[str, Any]]]:
        """
        Detects file extension and extracts text page-by-page.
        Returns a list of dicts with 'page' and 'text'.
        If text cannot be extracted or format is unsupported, returns None.
        """
        ext = Path(file_name).suffix.lower()
        try:
            if ext == ".pdf":
                return DocumentService._extract_pdf(file_bytes)
            elif ext == ".docx":
                return DocumentService._extract_docx(file_bytes)
            elif ext == ".xlsx":
                return DocumentService._extract_xlsx(file_bytes)
            elif ext == ".pptx":
                return DocumentService._extract_pptx(file_bytes)
            else:
                logger.warning(f"Unsupported file format: {ext}")
                return None
        except Exception as e:
            logger.error(f"Error during {ext} text extraction: {e}", exc_info=True)
            return None

    @staticmethod
    def _extract_pdf(pdf_bytes: bytes) -> Optional[List[Dict[str, Any]]]:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        pages_data = []
        total_text_length = 0
        
        for page_num, page in enumerate(doc):
            text = page.get_text()
            total_text_length += len(text.strip())
            pages_data.append({
                "page": page_num + 1,
                "text": text
            })
            
        if total_text_length == 0:
            return None
            
        return pages_data

    @staticmethod
    def _extract_docx(docx_bytes: bytes) -> Optional[List[Dict[str, Any]]]:
        from io import BytesIO
        doc = docx.Document(BytesIO(docx_bytes))
        pages_data = []
        
        # Word documents don't have natural "pages" in raw XML.
        # We group paragraphs (e.g. 5 paragraphs per logical "page") to support citations.
        current_page_paragraphs = []
        page_num = 1
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                current_page_paragraphs.append(text)
            
            if len(current_page_paragraphs) >= 5:
                pages_data.append({
                    "page": page_num,
                    "text": "\n".join(current_page_paragraphs)
                })
                current_page_paragraphs = []
                page_num += 1
                
        # Append remaining paragraphs
        if current_page_paragraphs:
            pages_data.append({
                "page": page_num,
                "text": "\n".join(current_page_paragraphs)
            })

        # Process tables
        for table_idx, table in enumerate(doc.tables):
            table_text_lines = []
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    table_text_lines.append(" | ".join(row_text))
            
            if table_text_lines:
                pages_data.append({
                    "page": page_num,
                    "text": f"[Table {table_idx + 1}]\n" + "\n".join(table_text_lines)
                })
                page_num += 1
                
        if not pages_data or sum(len(p["text"]) for p in pages_data) == 0:
            return None
            
        return pages_data

    @staticmethod
    def _extract_xlsx(xlsx_bytes: bytes) -> Optional[List[Dict[str, Any]]]:
        from io import BytesIO
        wb = openpyxl.load_workbook(BytesIO(xlsx_bytes), data_only=True)
        pages_data = []
        page_num = 1
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_rows = []
            
            for row in sheet.iter_rows(values_only=True):
                # Filter out completely empty rows
                row_vals = [str(val).strip() for val in row if val is not None]
                if row_vals:
                    sheet_rows.append(" | ".join(row_vals))
            
            if sheet_rows:
                # Treat each worksheet as a separate logical "page"
                pages_data.append({
                    "page": page_num,
                    "text": f"Sheet: {sheet_name}\n" + "\n".join(sheet_rows)
                })
                page_num += 1
                
        if not pages_data or sum(len(p["text"]) for p in pages_data) == 0:
            return None
            
        return pages_data

    @staticmethod
    def _extract_pptx(pptx_bytes: bytes) -> Optional[List[Dict[str, Any]]]:
        from io import BytesIO
        prs = pptx.Presentation(BytesIO(pptx_bytes))
        pages_data = []
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_text_elements = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text_elements.append(shape.text.strip())
            
            if slide_text_elements:
                # Treat each slide as a separate "page"
                pages_data.append({
                    "page": slide_idx + 1,
                    "text": f"[Slide {slide_idx + 1}]\n" + "\n".join(slide_text_elements)
                })
                
        if not pages_data or sum(len(p["text"]) for p in pages_data) == 0:
            return None
            
        return pages_data

    @staticmethod
    def save_extracted_text(document_id: str, pages_data: List[Dict[str, Any]]) -> str:
        """
        Saves the extracted pages data as a JSON file locally.
        Path: backend/uploads/processed/{documentId}.json
        """
        processed_dir = Path(__file__).parent.parent / "uploads" / "processed"
        processed_dir.mkdir(parents=True, exist_ok=True)
        
        output_path = processed_dir / f"{document_id}.json"
        
        document_json = {
            "documentId": document_id,
            "totalPages": len(pages_data),
            "pages": pages_data
        }
        
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(document_json, f, ensure_ascii=False, indent=2)
            
        return str(output_path)
